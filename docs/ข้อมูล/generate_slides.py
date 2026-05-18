"""
Generate PPTX presentation: คู่มือระบบประเมิน 360 องศา กนอ.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
CARD_BG = RGBColor(0x1E, 0x29, 0x3B)
BLUE = RGBColor(0x60, 0xA5, 0xFA)
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
PINK = RGBColor(0xF4, 0x72, 0xB6)
GREEN = RGBColor(0x34, 0xD3, 0x99)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
WHITE = RGBColor(0xF1, 0xF5, 0xF9)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)
DARK_BLUE = RGBColor(0x1E, 0x3A, 0x5F)
ROSE = RGBColor(0xFB, 0x71, 0x85)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=16, color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, space_before=0, space_after=0, font_name='Segoe UI'):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    return p


def add_rounded_rect(slide, left, top, width, height, fill_color=CARD_BG):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title, body, icon="",
             title_color=WHITE, fill_color=CARD_BG):
    rect = add_rounded_rect(slide, left, top, width, height, fill_color)
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)

    p = tf.paragraphs[0]
    p.text = f"{icon} {title}" if icon else title
    p.font.size = Pt(16)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.font.name = 'Segoe UI'

    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(13)
    p2.font.color.rgb = GRAY
    p2.font.name = 'Segoe UI'
    p2.space_before = Pt(6)
    return rect


def add_slide_number(slide, num, total=16):
    add_textbox(slide, 12.0, 7.0, 1.2, 0.4, f"{num}/{total}",
                font_size=11, color=GRAY, alignment=PP_ALIGN.RIGHT)


def add_table(slide, left, top, width, height, rows, cols, data,
              col_widths=None, header_color=DARK_BLUE, font_size=13):
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                          Inches(width), Inches(height))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = 'Segoe UI'
                if r == 0:
                    paragraph.font.color.rgb = BLUE
                    paragraph.font.bold = True
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = header_color
                else:
                    paragraph.font.color.rgb = LIGHT_GRAY
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = CARD_BG if r % 2 == 0 else DARK_BG
    return table_shape


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

add_textbox(slide, 1, 1.5, 11, 1.2,
            "ระบบประเมิน 360 องศา กนอ.",
            font_size=48, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1, 2.8, 11, 0.8,
            "คู่มือการใช้งานระบบประเมินผลบุคลากร",
            font_size=24, color=GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1, 3.5, 11, 0.6,
            "การนิคมอุตสาหกรรมแห่งประเทศไทย",
            font_size=20, color=PURPLE, alignment=PP_ALIGN.CENTER)

# URL box
rect = add_rounded_rect(slide, 4.5, 4.5, 4.3, 0.6, CARD_BG)
tf = rect.text_frame
tf.paragraphs[0].text = "evaluation.milesconsult.com"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.color.rgb = BLUE
tf.paragraphs[0].font.name = 'Segoe UI'
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

add_textbox(slide, 1, 5.8, 11, 0.5,
            "จัดทำโดย Miles Consult  |  ปรับปรุง 22 มีนาคม 2569",
            font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 1)

# ============================================================
# SLIDE 2: Overview - 5 angles
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.7, 0.4, 10, 0.7, "ภาพรวมระบบ — การประเมิน 5 มุมมอง",
            font_size=32, color=BLUE, bold=True)
add_textbox(slide, 0.7, 1.0, 10, 0.5,
            "ระบบประเมิน 360 องศาที่รองรับการประเมินจากหลายมุมมอง",
            font_size=16, color=GRAY)

cards = [
    ("Top-down (บนลงล่าง)", "ผู้บังคับบัญชาประเมินผู้ใต้บังคับบัญชา"),
    ("Bottom-up (ล่างขึ้นบน)", "ผู้ใต้บังคับบัญชาประเมินผู้บังคับบัญชา"),
    ("Left / Peer (ระดับเดียวกัน)", "เพื่อนร่วมงานระดับเดียวกันประเมินกัน"),
    ("Right / Cross (ข้ามสายงาน)", "ประเมินข้ามหน่วยงาน + ผู้ประเมินภายนอก"),
    ("Self (ประเมินตนเอง)", "พนักงานประเมินผลงานตนเอง"),
]
icons = ["arrow-down", "arrow-up", "arrows-lr", "refresh", "mirror"]

x_positions = [0.5, 2.9, 5.3, 7.7, 10.1]
for i, (title, body) in enumerate(cards):
    add_card(slide, x_positions[i], 1.8, 2.2, 1.6, title, body, fill_color=CARD_BG)

# Grade groups
add_textbox(slide, 0.7, 3.8, 10, 0.6, "กลุ่มผู้ถูกประเมิน", font_size=24, color=PURPLE, bold=True)

data = [
    ["กลุ่ม", "ระดับ (Grade)", "คำอธิบาย"],
    ["ผู้ว่าการ", "13", "ผู้บริหารสูงสุด"],
    ["ผู้บริหาร", "9 – 12", "ผู้บริหารระดับสูง–กลาง"],
    ["พนักงาน", "4 – 8", "พนักงานทั่วไป"],
]
add_table(slide, 0.7, 4.4, 6, 1.8, 4, 3, data, col_widths=[2, 2, 2])

# Fiscal year box
rect = add_rounded_rect(slide, 7.2, 4.4, 5.5, 1.8, CARD_BG)
tf = rect.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
tf.margin_top = Inches(0.15)
p = tf.paragraphs[0]
p.text = "ปีงบประมาณ (Fiscal Year)"
p.font.size = Pt(18)
p.font.color.rgb = AMBER
p.font.bold = True
p.font.name = 'Segoe UI'
add_paragraph(tf, "ระบบใช้ปีงบไทย: ตุลาคม – กันยายน", font_size=14, color=LIGHT_GRAY, space_before=8)
add_paragraph(tf, "ต.ค. 2568 – ก.ย. 2569 = ปีงบ 2569", font_size=14, color=GRAY, space_before=4)
add_paragraph(tf, "Admin เลือกดูปีงบต่างๆ ได้", font_size=14, color=GRAY, space_before=4)

add_slide_number(slide, 2)

# ============================================================
# SLIDE 3: Login
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.7, 0.4, 10, 0.7, "การเข้าสู่ระบบ",
            font_size=32, color=BLUE, bold=True)

# Internal login
add_textbox(slide, 0.7, 1.2, 5, 0.5, "บุคลากรภายใน — /login",
            font_size=22, color=PURPLE, bold=True)

data = [
    ["ฟิลด์", "รายละเอียด"],
    ["รหัสพนักงาน (emid)", "ตัวเลข 6 หลัก เช่น 350101"],
    ["รหัสผ่าน", "รหัสผ่านส่วนตัว"],
]
add_table(slide, 0.7, 1.8, 5.5, 1.2, 3, 2, data, col_widths=[2.5, 3])

# Example accounts
add_textbox(slide, 0.7, 3.2, 5, 0.5, "บัญชีตัวอย่าง (Seeder)",
            font_size=20, color=AMBER, bold=True)

data2 = [
    ["บทบาท", "รหัสพนักงาน", "ชื่อ", "รหัสผ่าน"],
    ["Admin", "350101", "ปคุณดา ชั้นบุญ", "13112541"],
    ["User", "000000", "อัฏฐพล นิติสุพรฯ", "13112541"],
]
add_table(slide, 0.7, 3.8, 6, 1.1, 3, 4, data2, col_widths=[1.2, 1.5, 2, 1.3])

# External login
add_textbox(slide, 7.2, 1.2, 5, 0.5, "ผู้ประเมินภายนอก — /external/login",
            font_size=22, color=PURPLE, bold=True)

data3 = [
    ["ฟิลด์", "รายละเอียด"],
    ["Access Code", "รหัสที่ Admin สร้างให้"],
]
add_table(slide, 7.2, 1.8, 5.5, 0.8, 2, 2, data3, col_widths=[2.5, 3])

# Notes
rect = add_rounded_rect(slide, 7.2, 2.8, 5.5, 2.0, CARD_BG)
tf = rect.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
tf.margin_top = Inches(0.15)
p = tf.paragraphs[0]
p.text = "คุณสมบัติ Access Code:"
p.font.size = Pt(16)
p.font.color.rgb = AMBER
p.font.bold = True
p.font.name = 'Segoe UI'
for t in ["- ใช้ได้ครั้งเดียว", "- มีวันหมดอายุ", "- ผูกกับองค์กรภายนอก", "- Admin สร้าง/เพิกถอน/สร้างใหม่ได้"]:
    add_paragraph(tf, t, font_size=14, color=LIGHT_GRAY, space_before=4)

# Warning
rect = add_rounded_rect(slide, 0.7, 5.2, 12, 0.8, DARK_BLUE)
tf = rect.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
p = tf.paragraphs[0]
p.text = "หมายเหตุ: ในระบบจริง (Production) ข้อมูลผู้ใช้จะถูกนำเข้าจากฐานข้อมูลจริงของ กนอ. ไม่ใช่ข้อมูล Seeder  |  ป้องกัน Brute-force: ล็อกหลัง 5 ครั้ง/นาที"
p.font.size = Pt(14)
p.font.color.rgb = AMBER
p.font.name = 'Segoe UI'

add_slide_number(slide, 3)

# ============================================================
# SLIDE 4: User Roles
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.7, 0.4, 10, 0.7, "บทบาทผู้ใช้งาน 3 ประเภท",
            font_size=32, color=BLUE, bold=True)

# User card
user_items = "- ประเมินตนเอง\n- ประเมินผู้อื่นตามมอบหมาย\n- ดูสถานะการประเมิน\n- ทำแบบสำรวจความพึงพอใจ\n- แก้ไขโปรไฟล์"
add_card(slide, 0.5, 1.3, 3.8, 3.5, "User (ผู้ใช้ทั่วไป)", user_items, title_color=BLUE)

# Admin card
admin_items = "- จัดการผู้ใช้งาน (CRUD)\n- จัดการโครงสร้างองค์กร\n- สร้าง/จัดการแบบประเมิน\n- มอบหมายการประเมิน\n- จัดการ Access Code\n- ดูรายงาน + Export"
add_card(slide, 4.7, 1.3, 3.8, 3.5, "Admin (ผู้ดูแลระบบ)", admin_items, title_color=PINK)

# External card
ext_items = "- เข้าระบบด้วย Access Code\n- ยืนยันตัวตนและองค์กร\n- ทำแบบประเมินตามที่กำหนด\n- ประเมินได้ครั้งเดียว/Code\n- ไม่ต้องสมัครสมาชิก"
add_card(slide, 8.9, 1.3, 3.8, 3.5, "External (ผู้ประเมินภายนอก)", ext_items, title_color=GREEN)

# Login redirect info
rect = add_rounded_rect(slide, 0.5, 5.2, 12.3, 1.2, DARK_BLUE)
tf = rect.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.1)
p = tf.paragraphs[0]
p.text = "การ Redirect หลัง Login"
p.font.size = Pt(18)
p.font.color.rgb = BLUE
p.font.bold = True
p.font.name = 'Segoe UI'
add_paragraph(tf, "Admin  →  /dashboardadmin  (แดชบอร์ดผู้ดูแล)     |     User  →  /dashboard  (แดชบอร์ดผู้ใช้)     |     External  →  /external/confirm  (ยืนยันตัวตน)", font_size=14, color=LIGHT_GRAY, space_before=6)

add_slide_number(slide, 4)

# ============================================================
# SLIDE 5: User Flow
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.7, 0.4, 10, 0.7, "ขั้นตอนการใช้งาน — ผู้ใช้ทั่วไป (User)",
            font_size=32, color=BLUE, bold=True)

# Step 1
add_textbox(slide, 0.7, 1.2, 10, 0.5, "1. เข้าสู่ระบบ → แดชบอร์ด",
            font_size=20, color=PURPLE, bold=True)
add_textbox(slide, 0.7, 1.7, 12, 0.4,
            "Login (/login)  →  Dashboard (/dashboard)  →  ดูสถานะการประเมินทั้งหมด  |  รายชื่อผู้ที่ต้องประเมิน  |  Progress bar",
            font_size=15, color=LIGHT_GRAY)

# Step 2
add_textbox(slide, 0.7, 2.3, 10, 0.5, "2. ประเมินตนเอง (Self-Evaluation)",
            font_size=20, color=PURPLE, bold=True)
add_textbox(slide, 0.7, 2.8, 12, 0.4,
            "เลือก 'ประเมินตนเอง'  →  ตอบทีละขั้น (Step-by-step)  →  Auto-save ทุกข้อ  →  กลับมาทำต่อได้ (Resume)  →  กดส่ง (Submit)",
            font_size=15, color=LIGHT_GRAY)

# Step 3
add_textbox(slide, 0.7, 3.4, 10, 0.5, "3. ประเมินผู้อื่น (Assigned Evaluation)",
            font_size=20, color=PURPLE, bold=True)
add_textbox(slide, 0.7, 3.9, 12, 0.4,
            "เลือกผู้ถูกประเมิน  →  ดูข้อมูล (ชื่อ/ตำแหน่ง/หน่วยงาน)  →  ตอบคำถาม  →  Peer Comparison Widget  →  กดส่ง",
            font_size=15, color=LIGHT_GRAY)

# Step 4
add_textbox(slide, 0.7, 4.5, 10, 0.5, "4. แบบสำรวจความพึงพอใจ (8 ข้อ, คะแนน 1-5)",
            font_size=20, color=PURPLE, bold=True)

data = [
    ["ข้อ", "หัวข้อ", "ข้อ", "หัวข้อ"],
    ["1", "ความพึงพอใจการใช้งาน", "5", "ความสะดวกในการเข้าถึง"],
    ["2", "ความง่ายในการใช้งาน", "6", "ความครบถ้วนของข้อมูล"],
    ["3", "ความเร็วในการตอบสนอง", "7", "ความเหมาะสมของเนื้อหา"],
    ["4", "ความถูกต้องของข้อมูล", "8", "ความพึงพอใจโดยรวม"],
]
add_table(slide, 0.7, 5.1, 12, 2.0, 5, 4, data, col_widths=[0.8, 4.5, 0.8, 4.5], font_size=12)

add_slide_number(slide, 5)

# ============================================================
# SLIDE 6: Question Types
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.7, 0.4, 10, 0.7, "ประเภทคำถามในแบบประเมิน",
            font_size=32, color=BLUE, bold=True)

data = [
    ["ประเภท", "คำอธิบาย", "ตัวอย่าง", "การบันทึก"],
    ["Rating", "เลือกคะแนน 1-5", "ดีเยี่ยม(5) / ดี(4) / พอใช้(3) / ปรับปรุง(2) / ปรับปรุงมาก(1)", "เก็บ option_id → score"],
    ["Open Text", "ตอบข้อความอิสระ", "จุดเด่นในการทำงาน / ข้อเสนอแนะ", "เก็บ text ใน value"],
    ["Choice", "เลือกตอบ 1 ข้อ", "เลือกระดับที่เหมาะสมที่สุด", "เก็บ option_id"],
    ["Multiple Choice", "เลือกได้หลายข้อ", "เลือกทักษะที่โดดเด่น (เลือกได้มากกว่า 1)", "เก็บ JSON array"],
]
add_table(slide, 0.7, 1.3, 12, 2.5, 5, 4, data, col_widths=[2, 3, 4, 3])

# Rating scale visual
add_textbox(slide, 0.7, 4.2, 10, 0.5, "ระดับคะแนน Rating", font_size=22, color=PURPLE, bold=True)

colors_rating = [
    (GREEN, "5 ดีเยี่ยม"),
    (BLUE, "4 ดี"),
    (PURPLE, "3 พอใช้"),
    (AMBER, "2 ปรับปรุง"),
    (ROSE, "1 ปรับปรุงมาก"),
]
for i, (c, label) in enumerate(colors_rating):
    rect = add_rounded_rect(slide, 0.7 + i * 2.5, 4.8, 2.2, 0.7, CARD_BG)
    tf = rect.text_frame
    tf.paragraphs[0].text = label
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = c
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = 'Segoe UI'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Auto-save note
rect = add_rounded_rect(slide, 0.7, 5.8, 12, 0.8, DARK_BLUE)
tf = rect.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
p = tf.paragraphs[0]
p.text = "ระบบ Auto-save ทุกข้อ  |  สามารถกลับมาทำต่อได้ (Resume)  |  ป้องกันการส่งซ้ำด้วย Unique Constraint"
p.font.size = Pt(15)
p.font.color.rgb = AMBER
p.font.name = 'Segoe UI'

add_slide_number(slide, 6)

# ============================================================
# SLIDE 7: Admin Dashboard
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.7, 0.4, 10, 0.7, "Admin Dashboard — /dashboardadmin",
            font_size=32, color=BLUE, bold=True)
add_textbox(slide, 0.7, 1.0, 10, 0.4,
            "ภาพรวมระบบสำหรับผู้ดูแล — เลือกปีงบประมาณได้",
            font_size=16, color=GRAY)

stats = [
    ("จำนวนผู้ใช้", "พนักงาน + Admin ทั้งหมด", BLUE),
    ("การมอบหมาย", "จำนวนทั้งหมด / ผู้ประเมิน / ผู้ถูกประเมิน", PURPLE),
    ("อัตราเสร็จสิ้น", "% ผู้ที่ทำประเมินเสร็จแล้ว", GREEN),
    ("แบบประเมิน", "จำนวนที่เผยแพร่แล้ว", AMBER),
    ("สถิติตามระดับ", "ผู้ว่าการ / ผู้บริหาร / พนักงาน", ROSE),
    ("สถิติตามมุมประเมิน", "top / bottom / left / right", BLUE),
    ("ผู้ประเมินภายนอก", "Access Code ทั้งหมด / ที่ใช้แล้ว", GREEN),
    ("ผลองค์กรภายนอก", "คะแนนเฉลี่ยแยกตามองค์กร", PURPLE),
]

for i, (title, desc, color) in enumerate(stats):
    row = i // 4
    col = i % 4
    x = 0.5 + col * 3.15
    y = 1.6 + row * 2.2
    rect = add_rounded_rect(slide, x, y, 2.95, 1.8, CARD_BG)
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(17)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = 'Segoe UI'
    add_paragraph(tf, desc, font_size=13, color=GRAY, space_before=8)

add_slide_number(slide, 7)

# ============================================================
# SLIDE 8: Admin - User & Org Management
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.7, 0.4, 10, 0.7, "Admin — จัดการผู้ใช้ & โครงสร้างองค์กร",
            font_size=32, color=BLUE, bold=True)

# User management
add_textbox(slide, 0.7, 1.2, 6, 0.5, "จัดการผู้ใช้ — /admin/users",
            font_size=20, color=PURPLE, bold=True)

data = [
    ["การดำเนินการ", "URL", "คำอธิบาย"],
    ["ดูรายชื่อ", "/admin/users", "ค้นหา, กรองตามหน่วยงาน/ระดับ"],
    ["เพิ่มผู้ใช้", "/admin/users/create", "กรอกรหัส, ชื่อ, สังกัด, ระดับ"],
    ["แก้ไขผู้ใช้", "/admin/users/{id}/edit", "แก้ไขข้อมูลทั้งหมดรวมบทบาท"],
    ["ลบผู้ใช้", "ปุ่มลบ", "ลบพร้อมข้อมูลที่เกี่ยวข้อง"],
]
add_table(slide, 0.7, 1.8, 6, 2.2, 5, 3, data, col_widths=[1.5, 2.3, 2.2], font_size=12)

# Org structure
add_textbox(slide, 7.2, 1.2, 5, 0.5, "โครงสร้างองค์กร",
            font_size=20, color=PURPLE, bold=True)

data2 = [
    ["ระดับ", "URL", "ภาษาไทย"],
    ["Division", "/admin/divisions", "สายงาน"],
    ["Department", "/admin/departments", "หน่วยงาน"],
    ["Position", "/admin/positions", "ตำแหน่ง"],
    ["Faction", "/admin/factions", "ฝ่าย (ข้ามสายงาน)"],
]
add_table(slide, 7.2, 1.8, 5.5, 2.2, 5, 3, data2, col_widths=[1.5, 2.3, 1.7], font_size=12)

# Hierarchy
add_textbox(slide, 0.7, 4.3, 5, 0.5, "ลำดับชั้นองค์กร",
            font_size=20, color=AMBER, bold=True)

rect = add_rounded_rect(slide, 0.7, 4.9, 5.5, 1.8, CARD_BG)
tf = rect.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.15)
p = tf.paragraphs[0]
p.text = "สายงาน (Division)"
p.font.size = Pt(16)
p.font.color.rgb = BLUE
p.font.bold = True
p.font.name = 'Consolas'
add_paragraph(tf, "  └── หน่วยงาน (Department)", font_size=15, color=PURPLE, font_name='Consolas', space_before=4)
add_paragraph(tf, "        └── ตำแหน่ง (Position)", font_size=15, color=GREEN, font_name='Consolas', space_before=4)
add_paragraph(tf, "", font_size=8, space_before=6)
add_paragraph(tf, "ฝ่าย (Faction) — อิสระจากลำดับชั้น", font_size=15, color=AMBER, font_name='Consolas', space_before=2)

# User fields
add_textbox(slide, 6.7, 4.3, 6, 0.5, "ฟิลด์ผู้ใช้ที่สำคัญ",
            font_size=20, color=AMBER, bold=True)

rect = add_rounded_rect(slide, 6.7, 4.9, 6, 1.8, CARD_BG)
tf = rect.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.1)
fields = [
    "emid — รหัสพนักงาน (ไม่ซ้ำ, 6 หลัก)",
    "prename / fname / lname — คำนำหน้า/ชื่อ/นามสกุล",
    "division / department / position / faction — สังกัด",
    "grade — ระดับ (4-13)",
    "role — บทบาท (user / admin)",
    "user_type — ประเภท (internal / external)",
]
for i, f in enumerate(fields):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = f"• {f}"
    p.font.size = Pt(12)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = 'Segoe UI'
    p.space_before = Pt(3)

add_slide_number(slide, 8)

# ============================================================
# SLIDE 9: Evaluation Management
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.7, 0.4, 10, 0.7, "Admin — จัดการแบบประเมิน",
            font_size=32, color=BLUE, bold=True)

# Structure
add_textbox(slide, 0.7, 1.1, 6, 0.5, "โครงสร้างแบบประเมิน",
            font_size=20, color=PURPLE, bold=True)

rect = add_rounded_rect(slide, 0.7, 1.6, 7, 3.0, CARD_BG)
tf = rect.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.15)
lines = [
    ("แบบประเมิน (Evaluation)", BLUE, "← /evaluations"),
    ("├── ส่วนที่ 1 (Part)", PURPLE, "← .../parts"),
    ("│   ├── ด้านที่ 1 (Aspect)", GREEN, "← .../aspects"),
    ("│   │   ├── คำถาม 1 (Rating 1-5)", LIGHT_GRAY, ""),
    ("│   │   └── คำถาม 2 (Open Text)", LIGHT_GRAY, ""),
    ("│   └── ด้านที่ 2 (มี Sub-aspects)", GREEN, ""),
    ("│       ├── ด้านย่อย 2.1", AMBER, "← .../subaspects"),
    ("│       │   └── คำถาม 3 (Choice)", LIGHT_GRAY, ""),
    ("│       └── ด้านย่อย 2.2", AMBER, ""),
    ("│           └── คำถาม 4 (Multiple Choice)", LIGHT_GRAY, ""),
    ("└── ส่วนที่ 2 (Part)", PURPLE, ""),
]
for i, (line, color, comment) in enumerate(lines):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    full = f"{line}  {comment}" if comment else line
    p.text = full
    p.font.size = Pt(12)
    p.font.color.rgb = color
    p.font.name = 'Consolas'
    p.space_before = Pt(1)

# Status flow
add_textbox(slide, 8.2, 1.1, 4, 0.5, "สถานะแบบประเมิน",
            font_size=20, color=PURPLE, bold=True)

statuses = [
    ("Draft", "สร้าง/แก้ไข", BLUE),
    ("Preview", "ดูตัวอย่าง", PURPLE),
    ("Published", "เผยแพร่ใช้งาน", GREEN),
    ("Inactive", "ปิดใช้งาน", AMBER),
]
for i, (name, desc, color) in enumerate(statuses):
    y = 1.7 + i * 0.7
    rect = add_rounded_rect(slide, 8.2, y, 4.5, 0.55, CARD_BG)
    tf = rect.text_frame
    tf.margin_left = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = f"{name}  —  {desc}"
    p.font.size = Pt(14)
    p.font.color.rgb = color
    p.font.name = 'Segoe UI'

# CRUD operations
add_textbox(slide, 0.7, 4.8, 10, 0.5, "การจัดการ", font_size=20, color=AMBER, bold=True)

data = [
    ["หน้า", "URL", "การดำเนินการ"],
    ["รายการแบบประเมิน", "/evaluations", "สร้าง / แก้ไข / ลบ / Preview / Publish"],
    ["จัดการส่วน (Part)", "/evaluations/{id}/parts", "สร้าง / แก้ไข / ลบ / เรียงลำดับ"],
    ["จัดการด้าน (Aspect)", ".../aspects", "สร้าง / แก้ไข / ลบ / กำหนด Sub-aspects"],
    ["จัดการคำถาม", ".../questions", "สร้าง / แก้ไข / ลบ / สร้าง Options + Score"],
]
add_table(slide, 0.7, 5.3, 12, 2.0, 5, 3, data, col_widths=[3, 3.5, 5.5], font_size=12)

add_slide_number(slide, 9)

# ============================================================
# SLIDE 10: Assignment Management
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.7, 0.4, 10, 0.7, "Admin — มอบหมายการประเมิน & Access Code",
            font_size=32, color=BLUE, bold=True)

# Assignments
add_textbox(slide, 0.7, 1.2, 6, 0.5, "มอบหมายการประเมิน — /admin/assignments",
            font_size=20, color=PURPLE, bold=True)

data = [
    ["การดำเนินการ", "คำอธิบาย"],
    ["สร้างรายบุคคล", "เลือก ผู้ประเมิน → ผู้ถูกประเมิน → มุม → ปีงบ"],
    ["สร้างจำนวนมาก (Bulk)", "สร้างหลายรายการพร้อมกัน"],
    ["แก้ไข / ลบ", "ทั้งรายบุคคลและ Bulk Delete"],
    ["Analytics", "สถิติภาพรวมการมอบหมาย"],
    ["Export", "ส่งออกข้อมูลการมอบหมาย"],
]
add_table(slide, 0.7, 1.8, 6, 2.5, 6, 2, data, col_widths=[2.5, 3.5], font_size=13)

# Angle table
add_textbox(slide, 7.2, 1.2, 5, 0.5, "มุมการประเมิน (Angle)",
            font_size=20, color=PURPLE, bold=True)

data2 = [
    ["มุม", "ความหมาย", "ตัวอย่าง"],
    ["top", "บนลงล่าง", "หัวหน้า → ลูกน้อง"],
    ["bottom", "ล่างขึ้นบน", "ลูกน้อง → หัวหน้า"],
    ["left", "ระดับเดียวกัน", "เพื่อนร่วมงาน"],
    ["right", "ข้ามสายงาน", "ต่างหน่วยงาน/ภายนอก"],
]
add_table(slide, 7.2, 1.8, 5.5, 2.2, 5, 3, data2, col_widths=[1, 2, 2.5], font_size=13)

# Access Code
add_textbox(slide, 0.7, 4.6, 6, 0.5, "จัดการ Access Code — /admin/access-codes",
            font_size=20, color=PURPLE, bold=True)

data3 = [
    ["การดำเนินการ", "คำอธิบาย"],
    ["Generate", "สร้าง Code — เลือกองค์กร, ผู้ถูกประเมิน, แบบประเมิน"],
    ["Revoke", "เพิกถอน Code"],
    ["Regenerate", "สร้าง Code ใหม่ทดแทน"],
    ["Export", "ส่งออกรายการ Code"],
    ["Print Cards", "พิมพ์บัตร Access Code"],
]
add_table(slide, 0.7, 5.2, 6, 2.2, 6, 2, data3, col_widths=[2, 4], font_size=12)

# External org
add_textbox(slide, 7.2, 4.6, 5, 0.5, "องค์กรภายนอก",
            font_size=20, color=PURPLE, bold=True)

rect = add_rounded_rect(slide, 7.2, 5.2, 5.5, 2.0, CARD_BG)
tf = rect.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
tf.margin_top = Inches(0.15)
p = tf.paragraphs[0]
p.text = "/admin/external-organizations"
p.font.size = Pt(14)
p.font.color.rgb = BLUE
p.font.name = 'Consolas'
for t in ["• สร้าง / แก้ไข / ลบองค์กร", "• กำหนดรหัสองค์กร (org_code)", "• เปิด/ปิดสถานะ (is_active)", "• ใช้ร่วมกับ Access Code"]:
    add_paragraph(tf, t, font_size=13, color=LIGHT_GRAY, space_before=5)

add_slide_number(slide, 10)

# ============================================================
# SLIDE 11: Reports
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.7, 0.4, 10, 0.7, "Admin — รายงานและ Export",
            font_size=32, color=BLUE, bold=True)
add_textbox(slide, 0.7, 1.0, 10, 0.4, "URL: /admin/reports/evaluation",
            font_size=16, color=GRAY)

data = [
    ["ประเภทรายงาน", "รายละเอียด", "รูปแบบ"],
    ["รายงานรวม ผู้บริหาร+พนักงาน", "ครบทั้ง Grade 9-12 และ 5-8", "Excel"],
    ["รายงานผู้บริหาร", "เฉพาะ Grade 9-12", "Excel"],
    ["รายงานพนักงาน", "เฉพาะ Grade 5-8", "Excel"],
    ["รายงานประเมินตนเอง", "ผลการประเมินตนเอง", "Excel"],
    ["รายงานผู้ว่าการ", "ระดับ 13 + Weighted Scoring", "Excel"],
    ["รายงานองค์กรภายนอก", "ผลจากผู้ประเมินภายนอก", "Excel"],
    ["รายงานรายบุคคล", "คะแนนเฉพาะบุคคล", "Excel / PDF"],
    ["รายงานเปรียบเทียบ", "เปรียบเทียบระหว่างบุคคล", "Excel"],
    ["ข้อมูลดิบ (Raw Data)", "คำตอบดิบทั้งหมด", "Excel"],
]
add_table(slide, 0.7, 1.5, 8, 4.0, 10, 3, data, col_widths=[3.5, 3, 1.5], font_size=12)

# Analytics tools
add_textbox(slide, 9.2, 1.5, 3.8, 0.5, "เครื่องมือวิเคราะห์",
            font_size=18, color=PURPLE, bold=True)

tools = [
    ("Individual Angle Report", "คะแนนรายบุคคลแยกมุม"),
    ("Assignment Tracking", "ติดตามสถานะมอบหมาย"),
    ("Completion Stats", "สถิติการทำเสร็จ"),
    ("Evaluatee Details", "รายละเอียดผู้ถูกประเมิน"),
    ("System Health", "สุขภาพระบบ"),
    ("Clear Cache", "ล้างแคชรายงาน"),
    ("Real-time Data", "ข้อมูลแบบ Real-time"),
    ("Dashboard Data API", "API สำหรับดึงข้อมูล"),
]

for i, (name, desc) in enumerate(tools):
    y = 2.1 + i * 0.55
    rect = add_rounded_rect(slide, 9.2, y, 3.8, 0.45, CARD_BG)
    tf = rect.text_frame
    tf.margin_left = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = f"{name}"
    p.font.size = Pt(11)
    p.font.color.rgb = GREEN
    p.font.bold = True
    p.font.name = 'Segoe UI'

add_slide_number(slide, 11)

# ============================================================
# SLIDE 12: External Flow
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.7, 0.4, 10, 0.7, "ขั้นตอน — ผู้ประเมินภายนอก (External)",
            font_size=32, color=BLUE, bold=True)

# Flow
steps = [
    ("1", "Admin สร้าง\nAccess Code", DARK_BLUE),
    ("2", "ส่ง Code ให้\nผู้ประเมินภายนอก", DARK_BLUE),
    ("3", "Login ด้วย Code\n/external/login", RGBColor(0x06, 0x4E, 0x3B)),
    ("4", "ยืนยันตัวตน\n/external/confirm", RGBColor(0x06, 0x4E, 0x3B)),
    ("5", "ทำแบบประเมิน\n/external/evaluate", RGBColor(0x06, 0x4E, 0x3B)),
    ("6", "ส่ง → ขอบคุณ\n/external/thank-you", RGBColor(0x06, 0x4E, 0x3B)),
]

for i, (num, text, bg) in enumerate(steps):
    x = 0.5 + i * 2.1
    rect = add_rounded_rect(slide, x, 1.3, 1.9, 1.4, bg)
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(28)
    p.font.color.rgb = BLUE
    p.font.bold = True
    p.font.name = 'Segoe UI'
    p.alignment = PP_ALIGN.CENTER
    add_paragraph(tf, text, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER, space_before=4)

# Screens table
add_textbox(slide, 0.7, 3.2, 10, 0.5, "หน้าจอผู้ประเมินภายนอก",
            font_size=22, color=PURPLE, bold=True)

data = [
    ["ลำดับ", "หน้า", "URL", "คำอธิบาย"],
    ["1", "เข้าสู่ระบบ", "/external/login", "กรอก Access Code"],
    ["2", "ยืนยันตัวตน", "/external/confirm", "ยืนยันชื่อและองค์กร"],
    ["3", "แดชบอร์ด", "/external/dashboard", "ดูข้อมูลผู้ถูกประเมิน"],
    ["4", "ทำแบบประเมิน", "/external/evaluate", "ตอบคำถามประเมิน"],
    ["5", "ขอบคุณ", "/external/thank-you", "ส่งสำเร็จ"],
]
add_table(slide, 0.7, 3.8, 12, 2.5, 6, 4, data, col_widths=[1, 2.5, 3.5, 5])

# Warning
rect = add_rounded_rect(slide, 0.7, 6.5, 12, 0.6, DARK_BLUE)
tf = rect.text_frame
tf.margin_left = Inches(0.3)
p = tf.paragraphs[0]
p.text = "Access Code ใช้ได้ครั้งเดียว  |  มีวันหมดอายุ  |  ต้องผูกกับองค์กรที่ active  |  Rate limit: 5 ครั้ง/นาที (login), 10 ครั้ง/นาที (submit)"
p.font.size = Pt(13)
p.font.color.rgb = AMBER
p.font.name = 'Segoe UI'

add_slide_number(slide, 12)

# ============================================================
# SLIDE 13: Weighted Scoring
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.7, 0.4, 10, 0.7, "การคำนวณคะแนน — Weighted Scoring (ผู้ว่าการ)",
            font_size=32, color=BLUE, bold=True)

# Stakeholder weights
add_textbox(slide, 0.7, 1.2, 6, 0.5, "น้ำหนักมุมผู้ประเมิน",
            font_size=20, color=PURPLE, bold=True)

data = [
    ["มุมประเมิน", "น้ำหนัก"],
    ["Top (บนลงล่าง)", "25%"],
    ["Bottom (ล่างขึ้นบน)", "25%"],
    ["Left (ระดับเดียวกัน)", "20%"],
    ["Right (ข้ามสายงาน)", "20%"],
    ["Self (ตนเอง)", "10%"],
]
add_table(slide, 0.7, 1.8, 5.5, 2.5, 6, 2, data, col_widths=[3.5, 2])

# Criteria weights
add_textbox(slide, 7.0, 1.2, 6, 0.5, "น้ำหนักเกณฑ์ประเมิน",
            font_size=20, color=PURPLE, bold=True)

data2 = [
    ["สมรรถนะ", "น้ำหนัก"],
    ["ภาวะผู้นำ (Leadership)", "30%"],
    ["วิสัยทัศน์ (Vision)", "25%"],
    ["การสื่อสาร (Communication)", "15%"],
    ["นวัตกรรม (Innovation)", "10%"],
    ["จริยธรรม (Ethics)", "10%"],
    ["ทำงานเป็นทีม (Teamwork)", "10%"],
]
add_table(slide, 7.0, 1.8, 5.5, 2.8, 7, 2, data2, col_widths=[3.5, 2])

# Calculation method
add_textbox(slide, 0.7, 4.8, 10, 0.5, "วิธีคำนวณ",
            font_size=20, color=AMBER, bold=True)

rect = add_rounded_rect(slide, 0.7, 5.3, 12, 1.5, CARD_BG)
tf = rect.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.15)

calc_steps = [
    ("1. คะแนนดิบ:", "คำนวณค่าเฉลี่ยของคะแนนทุกคำถามในแต่ละด้าน"),
    ("2. คะแนนถ่วงน้ำหนักด้าน:", "คูณค่าเฉลี่ยด้วยน้ำหนักของแต่ละด้าน (เช่น Leadership 30%)"),
    ("3. คะแนนถ่วงน้ำหนักมุม:", "รวมคะแนนจากทุกมุมตามน้ำหนัก (เช่น Top 25%, Bottom 25%)"),
    ("4. คะแนนรวม:", "ผลรวมคะแนนถ่วงน้ำหนักทั้งหมด → คะแนนสุดท้ายของผู้ถูกประเมิน"),
]
for i, (step, desc) in enumerate(calc_steps):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = f"{step} {desc}"
    p.font.size = Pt(13)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = 'Segoe UI'
    p.space_before = Pt(4)

add_slide_number(slide, 13)

# ============================================================
# SLIDE 14: Server Info
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.7, 0.4, 10, 0.7, "ข้อมูล Server และการดูแลระบบ",
            font_size=32, color=BLUE, bold=True)

# Server info
add_textbox(slide, 0.7, 1.2, 6, 0.5, "ข้อมูลเชื่อมต่อ",
            font_size=20, color=PURPLE, bold=True)

data = [
    ["รายการ", "ค่า"],
    ["Hosting", "Hostinger Cloud Enterprise"],
    ["Domain", "evaluation.milesconsult.com"],
    ["SSH IP", "156.67.209.215"],
    ["SSH Port", "65002"],
    ["Username", "u917560495"],
    ["SSH Auth", "SSH Key (~/.ssh/hostinger_deploy)"],
]
add_table(slide, 0.7, 1.8, 5.5, 2.8, 7, 2, data, col_widths=[1.5, 4])

# Database info
add_textbox(slide, 7.0, 1.2, 5, 0.5, "ข้อมูล Database",
            font_size=20, color=PURPLE, bold=True)

data2 = [
    ["รายการ", "ค่า"],
    ["Host", "localhost"],
    ["Database", "u917560495_milesconsultdb"],
    ["Username", "u917560495_jo_evaluation"],
    ["phpMyAdmin", "Hostinger Panel → DB"],
]
add_table(slide, 7.0, 1.8, 5.5, 2.0, 5, 2, data2, col_widths=[1.5, 4])

# SSH command
add_textbox(slide, 7.0, 4.0, 5, 0.5, "SSH Command",
            font_size=18, color=AMBER, bold=True)

rect = add_rounded_rect(slide, 7.0, 4.5, 5.5, 0.8, RGBColor(0x0D, 0x11, 0x17))
tf = rect.text_frame
tf.margin_left = Inches(0.2)
tf.margin_top = Inches(0.05)
p = tf.paragraphs[0]
p.text = "ssh -i ~/.ssh/hostinger_deploy \\\n  -p 65002 u917560495@156.67.209.215"
p.font.size = Pt(12)
p.font.color.rgb = GREEN
p.font.name = 'Consolas'

# Common commands
add_textbox(slide, 0.7, 4.8, 12, 0.5, "คำสั่งที่ใช้บ่อย",
            font_size=20, color=AMBER, bold=True)

rect = add_rounded_rect(slide, 0.7, 5.3, 5.8, 1.8, RGBColor(0x0D, 0x11, 0x17))
tf = rect.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
tf.margin_top = Inches(0.1)
cmds1 = [
    "# เข้าโปรเจค",
    "cd domains/evaluation.milesconsult.com/",
    "   laravel_project",
    "",
    "# ดู Log",
    "tail -50 storage/logs/laravel.log",
    "tail -f storage/logs/laravel.log",
]
for i, cmd in enumerate(cmds1):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = cmd
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY if cmd.startswith("#") else GREEN
    p.font.name = 'Consolas'

rect = add_rounded_rect(slide, 6.8, 5.3, 5.8, 1.8, RGBColor(0x0D, 0x11, 0x17))
tf = rect.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
tf.margin_top = Inches(0.1)
cmds2 = [
    "# ล้าง Cache",
    "php artisan cache:clear",
    "php artisan config:cache",
    "php artisan route:cache",
    "php artisan view:cache",
    "# ตรวจ Migration",
    "php artisan migrate:status",
]
for i, cmd in enumerate(cmds2):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = cmd
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY if cmd.startswith("#") else GREEN
    p.font.name = 'Consolas'

add_slide_number(slide, 14)

# ============================================================
# SLIDE 15: Troubleshooting
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.7, 0.4, 10, 0.7, "การแก้ไขปัญหาเบื้องต้น",
            font_size=32, color=BLUE, bold=True)

# Problem cards
problems = [
    ("500 Error", "tail -50 storage/logs/laravel.log\nchmod -R 755 storage bootstrap/cache\nphp artisan config:clear\nphp artisan cache:clear\nphp artisan config:cache", ROSE),
    ("หน้าขาว / ไม่มี CSS", "ls public_html/build/manifest.json\nrm -f public_html/hot\n\n→ ถ้าไม่มี manifest: ต้อง build + deploy ใหม่", AMBER),
    ("เข้าสู่ระบบไม่ได้", "• ตรวจรหัส 6 หลัก\n• ตรวจรหัสผ่าน\n• หากพยายาม > 5 ครั้ง → รอ 1 นาที\n• ติดต่อ Admin reset password", PURPLE),
    ("Access Code ไม่ทำงาน", "• ตรวจพิมพ์ถูกต้อง\n• Code ใช้ได้ครั้งเดียว\n• ตรวจวันหมดอายุ\n• ตรวจองค์กร active\n• ติดต่อ Admin regenerate", BLUE),
    ("Database Error", "php artisan migrate:status\nตรวจ .env → DB_HOST, DB_DATABASE\nphp artisan tinker → DB::connection()->getPdo()", GREEN),
    ("Composer Error", "rm -rf vendor\ncomposer install --no-dev\n  --optimize-autoloader\n  --no-interaction", AMBER),
]

for i, (title, body, color) in enumerate(problems):
    row = i // 3
    col = i % 3
    x = 0.5 + col * 4.15
    y = 1.2 + row * 3.0
    rect = add_rounded_rect(slide, x, y, 3.95, 2.7, CARD_BG)
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = 'Segoe UI'
    add_paragraph(tf, body, font_size=11, color=LIGHT_GRAY, font_name='Consolas', space_before=8)

add_slide_number(slide, 15)

# ============================================================
# SLIDE 16: Summary / Thank you
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 1, 0.8, 11, 0.8, "สรุป",
            font_size=44, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)

summaries = [
    ("ระบบประเมิน 360 องศา — 5 มุมประเมิน", GREEN),
    ("3 บทบาท: User, Admin, External", BLUE),
    ("โครงสร้างแบบประเมินแบบลำดับชั้น (Evaluation → Part → Aspect → Question)", PURPLE),
    ("Weighted Scoring สำหรับผู้ว่าการ (Grade 13)", AMBER),
    ("Export รายงาน Excel/PDF หลายรูปแบบ (10+ ประเภท)", GREEN),
    ("ระบบ Access Code สำหรับผู้ประเมินภายนอก", BLUE),
    ("รองรับหลายปีงบประมาณ (Fiscal Year Selector)", PURPLE),
    ("Hosted บน Hostinger Cloud Enterprise", AMBER),
]

for i, (text, color) in enumerate(summaries):
    y = 1.8 + i * 0.55
    rect = add_rounded_rect(slide, 1.5, y, 10, 0.45, CARD_BG)
    tf = rect.text_frame
    tf.margin_left = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = f"  {text}"
    p.font.size = Pt(16)
    p.font.color.rgb = color
    p.font.name = 'Segoe UI'

# URL
rect = add_rounded_rect(slide, 4.5, 6.0, 4.3, 0.5, CARD_BG)
tf = rect.text_frame
tf.paragraphs[0].text = "evaluation.milesconsult.com"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.color.rgb = BLUE
tf.paragraphs[0].font.name = 'Segoe UI'
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

add_textbox(slide, 1, 6.7, 11, 0.4,
            "Miles Consult  |  2569",
            font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 16)

# ============================================================
# SAVE
# ============================================================
output_path = r"C:\00_miles\Evaluation_project_miles_01-main\docs\user_manual_360_evaluation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
